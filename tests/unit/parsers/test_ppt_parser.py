import base64

import pytest
from scholar_lens.parsers.ppt_parser import PPTParser, _extract_shape_text


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)


def _make_pptx_with_speaker_notes(path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Visible Self-attention"
    slide.placeholders[1].text = "Visible Q K V explanation"
    slide.notes_slide.notes_text_frame.text = "Hidden teacher-only reminder"
    prs.save(path)


class TestPPTParser:
    def test_instantiation(self):
        parser = PPTParser()
        assert parser is not None
        assert hasattr(parser, "parse")

    def test_parse_nonexistent_raises(self):
        parser = PPTParser()
        with pytest.raises(FileNotFoundError):
            parser.parse("/nonexistent/file.pptx")

    def test_parse_excludes_speaker_notes_by_default(self, tmp_path):
        pptx_path = tmp_path / "slides.pptx"
        _make_pptx_with_speaker_notes(pptx_path)

        parsed = PPTParser().parse(pptx_path)

        assert "Visible Self-attention" in parsed.pages[0].text
        assert "Visible Q K V explanation" in parsed.raw_text
        assert "Hidden teacher-only reminder" not in parsed.pages[0].text
        assert "Hidden teacher-only reminder" not in parsed.raw_text
        assert "Hidden teacher-only reminder" not in parsed.sections[0]["text"]

    def test_pptx_chunks_exclude_speaker_notes(self, tmp_path):
        from scholar_lens.parsers.chunker import SectionAwareChunker

        pptx_path = tmp_path / "slides.pptx"
        _make_pptx_with_speaker_notes(pptx_path)

        parsed = PPTParser().parse(pptx_path)
        chunks = SectionAwareChunker().chunk(parsed, doc_id="deck")

        assert chunks
        assert any("Visible Self-attention" in chunk.text for chunk in chunks)
        assert all("Hidden teacher-only reminder" not in chunk.text for chunk in chunks)

    def test_extract_shape_text_reads_xml_without_text_frame_fallback(self):
        class FakeTextNode:
            def __init__(self, text):
                self.text = text

        class FakeParagraph:
            def __init__(self, texts):
                self._texts = texts

            def xpath(self, query):
                assert query == ".//a:t | .//m:t"
                return [FakeTextNode(text) for text in self._texts]

        class FakeElement:
            def xpath(self, query):
                if query == ".//a:p":
                    return [
                        FakeParagraph(["Self-attention ", "𝒒", "∙", "𝒌"]),
                        FakeParagraph(["Attention Matrix"]),
                    ]
                raise AssertionError(f"unexpected xpath query: {query}")

        class ExplodingTextFrame:
            @property
            def text(self):
                raise AssertionError("text_frame.text should not be used for PPTX math text")

        class FakeShape:
            _element = FakeElement()
            text_frame = ExplodingTextFrame()

        assert _extract_shape_text(FakeShape()) == "Self-attention 𝒒∙𝒌\nAttention Matrix"

    def test_parse_records_image_table_and_formula_metadata(self, tmp_path):
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "visual.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Visual attention slide"
        slide.shapes.add_picture(BytesIO(PNG_1X1), Inches(1), Inches(1), width=Inches(3), height=Inches(2))
        slide.shapes.add_table(2, 2, Inches(1), Inches(3.5), Inches(3), Inches(1))
        textbox = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = "Formula: softmax(QK^T / sqrt(d_k)) V"
        prs.save(pptx_path)

        parsed = PPTParser().parse(pptx_path)

        assert parsed.images
        assert parsed.images[0]["page"] == 0
        assert parsed.images[0]["area_ratio"] > 0
        assert parsed.tables
        assert parsed.tables[0]["page"] == 0
        assert parsed.formulas
        assert "softmax" in parsed.formulas[0]["text"]
