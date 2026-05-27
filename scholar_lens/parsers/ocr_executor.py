from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Any, Callable

from pydantic import BaseModel, Field

from scholar_lens.parsers.ocr_diagnostics import evaluate_ocr_result


class OCRUnavailableError(RuntimeError):
    pass


class OCRPageEnhancement(BaseModel):
    page: int
    text: str = ""
    ocr_quality: str = "failed"
    vision_recommended: bool = False
    reason: str = ""
    error: str = ""


class OCREnhancementResult(BaseModel):
    status: str = "completed"
    engine: str = "rapidocr"
    pages: list[OCRPageEnhancement] = Field(default_factory=list)
    error: str = ""


class RapidOCRExecutor:
    def __init__(
        self,
        ocr_callable: Callable[[str], Any] | None = None,
        prefer_gpu: bool = True,
    ) -> None:
        self._ocr_callable = ocr_callable
        self._prefer_gpu = prefer_gpu

    def run(self, source_path: str | Path, pages: list[int]) -> OCREnhancementResult:
        source = Path(source_path)
        suffix = source.suffix.lower()
        if suffix not in {".pdf", ".pptx"}:
            raise OCRUnavailableError("RapidOCR enhancement currently supports PDF and PPTX sources only")
        if not pages:
            return OCREnhancementResult(status="skipped", pages=[])

        ocr = self._ocr_callable or self._load_ocr_callable()
        page_results: list[OCRPageEnhancement] = []
        with tempfile.TemporaryDirectory(prefix="scholar_lens_ocr_") as tmp:
            work_dir = Path(tmp)
            rendered = (
                extract_pptx_slide_images(source, pages, work_dir)
                if suffix == ".pptx"
                else render_pdf_pages(source, pages, work_dir)
            )
            for page in pages:
                image_paths = rendered.get(page, [])
                if not image_paths:
                    reason = "pptx_no_embedded_images" if suffix == ".pptx" else "render_failed"
                    error = (
                        f"Slide {page} has no embedded images for lightweight PPTX OCR"
                        if suffix == ".pptx"
                        else f"Page {page} could not be rendered"
                    )
                    page_results.append(OCRPageEnhancement(
                        page=page,
                        ocr_quality="failed",
                        vision_recommended=True,
                        reason=reason,
                        error=error,
                    ))
                    continue
                try:
                    text = "\n".join(
                        normalize_rapidocr_output(ocr(str(image_path)))
                        for image_path in image_paths
                    ).strip()
                    quality = evaluate_ocr_result(text)
                    page_results.append(OCRPageEnhancement(
                        page=page,
                        text=quality.usable_text,
                        ocr_quality=quality.ocr_quality,
                        vision_recommended=quality.vision_recommended,
                        reason=quality.reason,
                    ))
                except Exception as exc:
                    page_results.append(OCRPageEnhancement(
                        page=page,
                        ocr_quality="failed",
                        vision_recommended=True,
                        reason="ocr_failed",
                        error=str(exc),
                    ))

        status = "completed"
        if page_results and all(page.ocr_quality == "failed" for page in page_results):
            status = "failed"
        return OCREnhancementResult(status=status, pages=page_results)

    def _load_ocr_callable(self) -> Callable[[str], Any]:
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError:
            try:
                from rapidocr import RapidOCR
            except ImportError as exc:
                raise OCRUnavailableError("RapidOCR is not installed") from exc

        engine = RapidOCR()
        return engine


def render_pdf_pages(source_path: Path, pages: list[int], work_dir: Path) -> dict[int, list[Path]]:
    import fitz

    rendered: dict[int, list[Path]] = {}
    doc = fitz.open(str(source_path))
    try:
        for page_num in pages:
            if page_num < 0 or page_num >= len(doc):
                continue
            page = doc[page_num]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = work_dir / f"page_{page_num}.png"
            pix.save(str(image_path))
            rendered[page_num] = [image_path]
    finally:
        doc.close()
    return rendered


def extract_pptx_slide_images(source_path: Path, pages: list[int], work_dir: Path) -> dict[int, list[Path]]:
    from pptx import Presentation

    rendered: dict[int, list[Path]] = {}
    prs = Presentation(str(source_path))
    for page_num in pages:
        if page_num < 0 or page_num >= len(prs.slides):
            continue
        slide = prs.slides[page_num]
        slide_dir = work_dir / f"slide_{page_num}"
        slide_dir.mkdir(parents=True, exist_ok=True)
        image_paths: list[Path] = []
        for idx, shape in enumerate(slide.shapes):
            image = getattr(shape, "image", None)
            if image is None:
                continue
            ext = (getattr(image, "ext", "") or "png").lower()
            image_path = slide_dir / f"image_{idx}.{ext}"
            image_path.write_bytes(image.blob)
            image_paths.append(image_path)
        if image_paths:
            rendered[page_num] = image_paths
    return rendered


def normalize_rapidocr_output(output: Any) -> str:
    lines: list[str] = []
    _collect_text_lines(output, lines)
    return "\n".join(line for line in lines if line)


def _collect_text_lines(value: Any, lines: list[str]) -> None:
    if isinstance(value, str):
        if value.strip():
            lines.append(value.strip())
        return
    if not isinstance(value, (list, tuple)):
        return
    if _looks_like_ocr_line(value):
        text = value[1]
        if isinstance(text, str) and text.strip():
            lines.append(text.strip())
        return
    for item in value:
        _collect_text_lines(item, lines)


def _looks_like_ocr_line(value: list | tuple) -> bool:
    return (
        len(value) >= 2
        and isinstance(value[1], str)
        and not isinstance(value[0], str)
    )
