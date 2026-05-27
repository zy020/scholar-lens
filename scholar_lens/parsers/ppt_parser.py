from __future__ import annotations

import re
from pathlib import Path

from scholar_lens.parsers.models import ParsedDocument, ParsedPage


_FORMULA_PATTERN = re.compile(
    r"(?:softmax|LayerNorm|sqrt|frac|sum|prod|argmax|argmin|\\[a-zA-Z]+|[A-Za-z]\s*[\^_=]\s*[A-Za-z0-9({])",
    re.IGNORECASE,
)


def _extract_shape_text(shape) -> str:
    """Extract PPTX text directly from XML nodes to avoid python-pptx math diagnostics."""
    element = getattr(shape, "_element", None)
    if element is None:
        return shape.text_frame.text.strip()

    paragraphs = []
    try:
        paragraph_nodes = element.xpath(".//a:p")
        for paragraph in paragraph_nodes:
            parts = [
                node.text or ""
                for node in paragraph.xpath(".//a:t | .//m:t")
                if node.text
            ]
            text = "".join(parts).strip()
            if text:
                paragraphs.append(text)
    except Exception:
        return shape.text_frame.text.strip()

    if paragraphs:
        return "\n".join(paragraphs)

    try:
        parts = [
            node.text or ""
            for node in element.xpath(".//a:t | .//m:t")
            if node.text
        ]
    except Exception:
        return shape.text_frame.text.strip()
    return "".join(parts).strip()


class PPTParser:
    """Parses PPTX files using python-pptx."""

    def parse(self, ppt_path: str | Path) -> ParsedDocument:
        ppt_path = Path(ppt_path)
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPTX not found: {ppt_path}")

        return self._parse_with_pptx(ppt_path)

    def _parse_with_pptx(self, ppt_path: Path) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(str(ppt_path))
        pages = []
        sections = []
        images: list[dict] = []
        tables: list[dict] = []
        formulas: list[dict] = []
        slide_area = max(float(prs.slide_width * prs.slide_height), 1.0)

        for i, slide in enumerate(prs.slides):
            texts = []
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = _extract_shape_text(shape)
                    if shape_text:
                        texts.append(shape_text)
                        if not title and shape.shape_type == 14:
                            title = shape_text
                        formulas.extend(_extract_formula_candidates(shape_text, i))
                if getattr(shape, "has_table", False):
                    tables.append(_shape_record(shape, i, "table", slide_area, len(tables)))
                if getattr(shape, "has_chart", False):
                    tables.append(_shape_record(shape, i, "chart", slide_area, len(tables)))
                if getattr(shape, "image", None) is not None:
                    images.append(_shape_record(shape, i, "image", slide_area, len(images)))

            slide_text = "\n".join(texts)

            pages.append(ParsedPage(page_num=i, text=slide_text, char_count=len(slide_text)))

            if title:
                sections.append({
                    "id": f"slide_{i}",
                    "title": title,
                    "level": 1,
                    "text": slide_text,
                    "page_start": i,
                    "page_end": i,
                })

        raw_text = "\n\n".join(p.text for p in pages)
        return ParsedDocument(
            source_path=str(ppt_path),
            doc_subtype="courseware_pptx",
            pages=pages,
            sections=sections,
            raw_text=raw_text,
            formulas=formulas,
            tables=tables,
            images=images,
        )


def _shape_record(shape, page_num: int, kind: str, slide_area: float, index: int) -> dict:
    left = float(getattr(shape, "left", 0) or 0)
    top = float(getattr(shape, "top", 0) or 0)
    width = float(getattr(shape, "width", 0) or 0)
    height = float(getattr(shape, "height", 0) or 0)
    return {
        "page": page_num,
        "index": index,
        "kind": kind,
        "bbox": [round(left, 2), round(top, 2), round(left + width, 2), round(top + height, 2)],
        "area_ratio": round(max(width * height, 0.0) / slide_area, 4),
        "source": "python-pptx",
    }


def _extract_formula_candidates(text: str, page_num: int) -> list[dict]:
    candidates = []
    for line in text.splitlines():
        normalized = line.strip()
        if len(normalized) < 4:
            continue
        if _FORMULA_PATTERN.search(normalized):
            candidates.append({
                "page": page_num,
                "text": normalized[:300],
                "source": "text_heuristic",
            })
    return candidates
